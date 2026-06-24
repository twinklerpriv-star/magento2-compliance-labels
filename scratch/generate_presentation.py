import sys
import os

# Configure console encoding to avoid issues with Unicode on Windows
sys.stdout.reconfigure(encoding='utf-8')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is not installed. Please install it using 'pip install python-pptx'.")
    sys.exit(1)

def main():
    dest_path = r"C:\Users\thomas.winkler\Desktop\Projekte\Google Antigravity\ELEKTROPEPI\wko_seminar_verbrauchsaenderungsgesetz.pptx"
    os.makedirs(os.path.dirname(dest_path), exist_ok=True)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide
    
    footer_text = "Elektro Pepi GmbH  |  Lieferanten-Information: Verbraucherrecht 2026"
    
    # Helper to add standard slide
    def add_slide(title, bullets=None, layout="standard", custom_content=None):
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = Inches(0)
        tf_title.margin_top = Inches(0)
        tf_title.margin_right = Inches(0)
        tf_title.margin_bottom = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(0, 0, 0)
        
        # 2. Slide Content
        if layout == "standard" and bullets:
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = Inches(0)
            tf_content.margin_top = Inches(0)
            tf_content.margin_right = Inches(0)
            tf_content.margin_bottom = Inches(0)
            
            for idx, bullet in enumerate(bullets):
                p = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
                
                # Determine indentation level
                level = 0
                clean_bullet = bullet
                if bullet.startswith("    - "):
                    level = 2
                    clean_bullet = bullet[6:]
                elif bullet.startswith("  - "):
                    level = 1
                    clean_bullet = bullet[4:]
                elif bullet.startswith("- "):
                    level = 0
                    clean_bullet = bullet[2:]
                
                p.text = clean_bullet
                p.level = level
                p.font.name = "Arial"
                p.font.size = Pt(18 - level * 2)
                p.font.color.rgb = RGBColor(30, 30, 30)
                p.space_after = Pt(8)
                
        elif layout == "bordered_box" and custom_content:
            # Draw bordered box for label simulation
            box_left = Inches(1.5)
            box_top = Inches(1.6)
            box_width = Inches(10.33)
            box_height = Inches(4.7)
            
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, box_top, box_width, box_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.line.width = Pt(1.5)
            
            tf_box = shape.text_frame
            tf_box.word_wrap = True
            tf_box.margin_left = Inches(0.4)
            tf_box.margin_right = Inches(0.4)
            tf_box.margin_top = Inches(0.3)
            tf_box.margin_bottom = Inches(0.3)
            
            # Header Inside Box
            p_box_title = tf_box.paragraphs[0]
            p_box_title.text = custom_content[0].upper()
            p_box_title.font.name = "Arial"
            p_box_title.font.size = Pt(16)
            p_box_title.font.bold = True
            p_box_title.font.color.rgb = RGBColor(0, 0, 0)
            p_box_title.alignment = PP_ALIGN.CENTER
            p_box_title.space_after = Pt(12)
            
            for line in custom_content[1:]:
                p = tf_box.add_paragraph()
                
                # Check formatting
                if line.startswith("  • ") or line.startswith("  1. ") or line.startswith("  2. "):
                    p.text = line[4:]
                    p.level = 1
                else:
                    p.text = line
                    p.level = 0
                    
                p.font.name = "Arial"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(30, 30, 30)
                p.space_after = Pt(4)
                
        elif layout == "translations" and custom_content:
            # 3 Columns for translations
            col_width = Inches(3.7)
            col_gap = Inches(0.3)
            left_margin = Inches(0.8)
            top_pos = Inches(2.2)
            height = Inches(4.2)
            
            num_items = len(custom_content)
            col_size = (num_items + 2) // 3
            
            for col_idx in range(3):
                start = col_idx * col_size
                end = min(start + col_size, num_items)
                col_items = custom_content[start:end]
                
                col_left = left_margin + col_idx * (col_width + col_gap)
                col_box = slide.shapes.add_textbox(col_left, top_pos, col_width, height)
                tf_col = col_box.text_frame
                tf_col.word_wrap = True
                tf_col.margin_left = Inches(0)
                tf_col.margin_top = Inches(0)
                tf_col.margin_right = Inches(0)
                tf_col.margin_bottom = Inches(0)
                
                for idx, item in enumerate(col_items):
                    p = tf_col.add_paragraph() if idx > 0 else tf_col.paragraphs[0]
                    p.text = item
                    p.font.name = "Arial"
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(30, 30, 30)
                    p.space_after = Pt(4)
                    
        # 3. Slide Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.3))
        tf_footer = footer_box.text_frame
        tf_footer.margin_left = Inches(0)
        tf_footer.margin_top = Inches(0)
        tf_footer.margin_right = Inches(0)
        tf_footer.margin_bottom = Inches(0)
        p_foot = tf_footer.paragraphs[0]
        p_foot.text = footer_text
        p_foot.font.name = "Arial"
        p_foot.font.size = Pt(9)
        p_foot.font.color.rgb = RGBColor(120, 120, 120)
        p_foot.alignment = PP_ALIGN.RIGHT
        
    # --- SLIDES DEFINITION ---
    
    # Slide 1: Welcome / Overview (Widerrufsbutton removed)
    add_slide(
        "Verbraucherrechts-Änderungsgesetz 2026 – Kurzüberblick",
        [
            "- Neue Informations- und Transparenzvorgaben im Online-Handel und im stationären Handel im Zusammenhang mit Gewährleistung und Garantie (Verwendung einer „harmonisierten Mitteilung“).",
            "- Vorgaben zur standardisierten Gestaltung von Nachhaltigkeitsinformationen: z. B. Informationen vor Vertragsabschluss über Reparierbarkeit, Ersatzteilverfügbarkeit und umweltfreundliche Lieferoptionen.",
            "- [Neue Regelungen für Fernabsatzverträge bei Finanzdienstleistungen: neue umfassende Informationspflichten & Vereinheitlichung des Rechtsrahmens]."
        ]
    )
    
    # Slide 2: Gewährleistungs- und Garantie-Label (Intro)
    add_slide(
        "Verbraucherrechts-Änderungsgesetz 2026 – Gewährleistungs- & Garantie-Label",
        [
            "- Neu: Verpflichtende Gewährleistungs- und Garantie-Label, welche Informations-, Kennzeichnungs- und Hinweispflichten für den stationären Handel als auch für Online-Shops mit sich bringen.",
            "- Die Pflichten zur Verwendung der Gewährleistungs- und Garantie-Label gelten für alle Händler, die an Verbraucher (B2C) verkaufen.",
            "- Die neuen Pflichten gelten für alle beweglichen körperlichen Waren, einschließlich solcher mit digitalen Elementen (z. B. Smartphones, vernetzte Küchengeräte).",
            "- Auch gebrauchte Waren fallen unter diese Regelung.",
            "- Ausgenommen sind reine B2B-Geschäfte."
        ]
    )
    
    # Slide 3: Gewährleistungs-Label (Kontext)
    add_slide(
        "Verbraucherrechts-Änderungsgesetz 2026 – Gewährleistungs-Label",
        [
            "- Hinweis auf Bestehen der gesetzlichen Gewährleistungspflicht (wie bisher).",
            "- Neu ist die verpflichtende Verwendung der europaweit „harmonisierten Mitteilung“ beim Verkauf von Waren.",
            "- Die harmonisierte Mitteilung ist neutral formuliert und darf inhaltlich nicht verändert werden.",
            "- Sie muss sowohl online im Webshop als auch im stationären Handel verwendet werden.",
            "- Online-Vorgabe: Das Label muss zwingend in Farbe angezeigt werden.",
            "- Stationäre Vorgabe: Im stationären Handel ist auch ein Schwarz-Weiß-Druck zulässig."
        ]
    )
    
    # Slide 4: Gewährleistungs-Label (Harmonisierte Mitteilung)
    gewaehrleistung_box = [
        "GESETZLICHE GEWÄHRLEISTUNG",
        "Mindestens zwei Jahre gesetzliche Gewährleistung der Vertragsmäßigkeit für Waren, die in der Europäischen Union verkauft werden.",
        "Verbraucherinnen und Verbraucher können ihre Rechte im Rahmen des gesetzlichen Gewährleistungsrechts geltend machen, z. B. wenn die Waren:",
        "  • nicht der Beschreibung entsprechen,",
        "  • nicht bestimmungsgemäß funktionieren.",
        "Verkäufer haften für jede Vertragswidrigkeit bei Lieferung, die innerhalb des Zeitraums der gesetzlichen Gewährleistung erkennbar wird. Verkäufer müssen in solchen Fällen anbieten:",
        "  • kostenlose Nachbesserung oder kostenlose Ersatzlieferung,",
        "  • in bestimmten Fällen eine Preisminderung oder eine vollständige Erstattung des Kaufpreises.",
        "In einigen Ländern gilt ein längerer Zeitraum für die gesetzliche Gewährleistung. Für gebrauchte Waren kann ein kürzerer Zeitraum gelten, jedoch nicht weniger als ein Jahr. Für weitere Informationen scannen Sie den QR-Code (europa.eu/youreurope/garantien) oder fragen Sie den Verkäufer.",
        "Was ist zu tun, wenn Sie vertragswidrige Waren erhalten?",
        "  1. Melden Sie dem Verkäufer das Problem so bald wie möglich.",
        "  2. Legen Sie einen Kaufnachweis vor, z. B. die Quittung, Rechnung oder einen Kontoauszug.",
        "Verkäufer und Hersteller können auch gewerbliche Garantien gewähren, die unabhängig von der gesetzlichen Gewährleistung gelten. Diese GARAN-Kennzeichnung zeigt beispielsweise, dass der Hersteller eine gewerbliche Haltbarkeitsgarantie ohne zusätzliche Kosten gewährt, die die gesamte Ware abdeckt."
    ]
    add_slide(
        "Gesetzliches Gewährleistungs-Label (Harmonisierte Mitteilung)",
        layout="bordered_box",
        custom_content=gewaehrleistung_box
    )
    
    # Slide 5: Gewährleistungs-Label - Technische & Stationäre Vorgaben
    add_slide(
        "Gewährleistungs-Label – Technische & Stationäre Vorgaben",
        [
            "- **Technische Vorgaben:** Technische Vorgaben wie Mindestgrößen, Schriftgrößen oder Schriftarten müssen bei der Darstellung im Webshop zwingend eingehalten werden.",
            "- **Darstellung im stationären Handel:**",
            "  - Die harmonisierte Mitteilung muss nicht an jedem einzelnen Produkt angebracht sein.",
            "  - Sie kann an einem gut sichtbaren, leicht wahrnehmbaren Ort platziert werden (z. B. als Plakat neben der Kasse oder im Eingangsbereich).",
            "- **QR-Code:** Der auf dem Label enthaltene QR-Code führt Verbraucher direkt zu einer EU-weit einheitlichen Informationsseite (mit konkreten Rechten und detaillierten Informationen)."
        ]
    )
    
    # Slide 6: Garantie-Label (Kontext)
    add_slide(
        "Verbraucherrechts-Änderungsgesetz 2026 – Garantie-Label",
        [
            "- Das harmonisierte Garantie-Label muss verwendet werden, wenn ein Hersteller dem Verbraucher eine kostenlose gewerbliche „Haltbarkeitsgarantie“ gewährt:",
            "  - diese muss für die gesamte Ware gelten (d. h. nicht nur für bestimmte Bestandteile oder Aspekte)",
            "  - sie muss eine Dauer von mehr als 2 Jahren umfassen UND",
            "  - der Hersteller stellt diese Garantieinformationen dem Händler zur Verfügung.",
            "- Wird in der Regel vom Hersteller auf der Verpackung angebracht.",
            "- Das Label muss immer entsprechend der gewährten Garantiedauer des jeweiligen Produktes angepasst werden („XX wird durch die Garantie-Jahre ersetzt“).",
            "- Die Angabe des Herstellers („Brand/Trademark“) und des konkreten Produktes („Model identifier“) müssen ebenfalls individuell angepasst werden."
        ]
    )
    
    # Slide 7: Garantie-Kennzeichnung (Standardisiertes GARAN-Label)
    add_slide(
        "Garantie-Kennzeichnung (Standardisiertes GARAN-Label)",
        [
            "- Das EU-Garantielabel (Kennzeichnung 'GARAN') ist visuell streng vorgegeben:",
            "  - Kopfzeile: Enthält den Schriftzug 'GARAN' mit einem Häkchen-Symbol und rechts das EU-G-Wappenschild mit Sternen.",
            "  - Produktzuordnung: Zeigt links die Marke ('Brand/Trademark') und rechts die Modellkennung ('Model identifier').",
            "  - Kern-Anzeige: Zeigt die Garantie-Jahre (z. B. '3' oder '5') in sehr großen Ziffern neben einem Kalender-Icon mit '365'.",
            "  - QR-Code: Verweist auf die offizielle EU-Informationsseite für gewerbliche Garantien.",
            "  - Fußzeile: Enthält den standardisierten mehrsprachigen Übersetzungstext für alle 24 EU-Amtssprachen.",
            "- Online-Shop-Darstellung:",
            "  - Das Label muss gut sichtbar auf der Produktdetailseite (PDP) integriert werden.",
            "  - Online kann auch eine platzsparende Kurzversion (horizontale Leiste) verwendet werden."
        ]
    )
    
    # Slide 8: Garantie-Label – Sprachen und Übersetzungen
    translations = [
        "BG: Гаранция от производителя в години",
        "CS: Záruka výrobce v letech",
        "DA: Producentgarantiens varighed i år",
        "DE: Herstellergarantie in Jahren",
        "EL: Eγγύηση παραγωγού σε έty",
        "EN: Producer guarantee in years",
        "ES: Garantía del productor en años",
        "ET: Tootja garantii aastates",
        "FI: Tuottajan takuu vuosina",
        "FR: Garantie du producteur en années",
        "GA: Ráthaíocht an táirgeora de réir blianta",
        "HR: Jamstvo proizvođača u godinama",
        "HU: Gyártói jótállás években",
        "IT: Garanzia del produttore in anni",
        "LT: Gamintojo garantija metais",
        "LV: Ražotāja garantija gados",
        "MT: Garanzija tal-produttur fi snin",
        "NL: Producentengarantie in jaren",
        "PL: Gwarancja producenta w latach",
        "PT: Garantia do produtor em anos",
        "RO: Garanția producătorului în ani",
        "SK: Záruka výrobcu v rokoch",
        "SL: Garancija proizvajalca v letih",
        "SV: Tillverkarens garanti i antal år"
    ]
    add_slide(
        "Garantie-Label – Sprachen und Übersetzungen",
        layout="translations",
        custom_content=translations
    )
    
    # Slide 9: Garantie-Label – Layout-Regeln & Formate
    add_slide(
        "Garantie-Label – Layout-Regeln & Formate",
        [
            "- **Inhaltliche Integrität:** Der restliche Inhalt des Labels (inkl. Farben, Logos und Anordnung) darf nicht verändert werden.",
            "- **Angabe der Dauer:** Die Dauer der Haltbarkeitsgarantie muss zwingend in vollen Jahren angegeben werden.",
            "  - Ausnahme: In Ausnahmefällen ist eine Angabe mit halben Jahren (z. B. '2,5' oder '3,5') zulässig.",
            "  - Andere Dezimalzahlen sind strikt verboten.",
            "- **Produktbezug:** Der Hinweis zur Garantie muss direkt beim jeweils betroffenen Produkt im Online-Shop oder im stationären Handel ersichtlich gemacht werden.",
            "- **Kurzversion:** Online kann alternativ die platzsparende Kurzversion verwendet werden:",
            "  - Visualisierung: [ XX 365 ] | GARAN (mit Häkchen) | EU-Schild.",
            "- **Farbdarstellung:** Online muss das Label immer in Farbe angezeigt werden. Im stationären Handel ist ein Schwarz-Weiß-Druck zulässig."
        ]
    )
    
    # Slide 10: Warenreparaturrichtlinie-Umsetzungsgesetz – Allgemeines
    add_slide(
        "Warenreparaturrichtlinie-Umsetzungsgesetz – Allgemeines",
        [
            "- **Geplantes Inkrafttreten:** Voraussichtlich am **31. Juli 2026** (soll für Verträge gelten, die ab diesem Stichtag abgeschlossen werden).",
            "  - Status: Befindet sich derzeit in Umsetzung; der formale Beschluss im Nationalrat steht noch aus.",
            "- **Gesetzliche Basis:** EU-Richtlinie zur Förderung der Reparatur von Waren (RL (EU) 2024/1799) – auch bekannt als „Recht auf Reparatur“-Richtlinie.",
            "- **Übergeordnetes Ziel:** Reparatur und Wiederverwendung brauchbarer Waren sowohl innerhalb als auch außerhalb der gesetzlichen Gewährleistung fördern.",
            "  - Verlängerung der Produktlebensdauer durch rechtliche und finanzielle Anreize.",
            "  - Erleichterung des Zugangs zu professionellen Reparaturdienstleistungen.",
            "- **Gesetzesänderungen:** Führt zu wesentlichen Anpassungen im Verbrauchergewährleistungsgesetz (VGG) und Konsumentenschutzgesetz (KSchG)."
        ]
    )
    
    # Slide 11: Warenreparaturrichtlinie – Wesentliche Inhalte
    add_slide(
        "Warenreparaturrichtlinie – Wesentliche Inhalte",
        [
            "- **Herstellerverpflichtung zur Reparatur:**",
            "  - Recht auf Reparatur für bestimmte Produktgruppen (Herstellerpflicht im B2C-Bereich).",
            "  - Diese Reparaturpflicht ist jedoch kostenpflichtig (der Verbraucher trägt die Kosten).",
            "- **Neue Informationspflichten:** Händler und Hersteller müssen transparente Informationen über verfügbare Reparaturdienste bereitstellen.",
            "- **Standardisiertes Formular:** Einführung eines standardisierten europäischen Formulars für Reparaturinformationen (zur besseren Vergleichbarkeit von Angeboten).",
            "- **Anpassungen im Gewährleistungsrecht:** Verlängerung der Gewährleistungsfrist unter bestimmten Bedingungen, wenn sich der Kunde für eine Reparatur statt eines Austauschs entscheidet."
        ]
    )
    
    # Slide 12: Auswirkungen in der Praxis – To-Dos für Lieferanten (Teil 1) (Widerrufsbutton removed)
    add_slide(
        "Auswirkungen in der Praxis – To-Dos für Lieferanten (Teil 1)",
        [
            "- **Anpassung der Kundenbeziehungen:** Verträge, AGB und Gewährleistungszusagen müssen rechtlich überprüft und an die neuen Standards angepasst werden.",
            "- **Bereitstellung von Garantiedaten:**",
            "  - Lieferanten müssen Händlern (wie Elektro Pepi) die Garantiedaten je Produkt strukturiert zur Verfügung stellen.",
            "  - Erforderliche Daten: Garantiedauer (in Jahren), Garantiebedingungen und Identität des Garantiegebers.",
            "  - Genaue Zuweisung von Marke ('Brand/Trademark') und Modellkennung ('Model identifier') je Artikel zur Schnittstellen-Übermittlung."
        ]
    )
    
    # Slide 13: Auswirkungen in der Praxis – To-Dos für Lieferanten (Teil 2)
    add_slide(
        "Auswirkungen in der Praxis – To-Dos für Lieferanten (Teil 2)",
        [
            "- **Klärung der Verpackungslogik:**",
            "  - Prüfung, ob die Garantielabels direkt vom Hersteller auf der Produktverpackung angebracht werden.",
            "  - Abklärung, ob SBS/Händler die Labels bei Blisterverpackungen im Lager selbst aufbringen müssen.",
            "- **Reparatur-Vereinbarungen:** Teilweise Anpassung der Verträge zwischen Herstellern und Händlern bezüglich der Reparaturabwicklung im Servicefall.",
            "- **Mitarbeiterschulung:** Durchführung von Schulungen für Einkauf, Vertrieb und Service bezüglich der neuen Gewährleistungs- und Garantieregelungen."
        ]
    )
    
    # Slide 14: Kontakt & Referenten-Info
    add_slide(
        "Kontakt & Referenten-Info",
        [
            "- **Referent des Seminars:**",
            "  - MMag. Stefan Adametz, LL.M., MBA (Rechtsanwalt)",
            "  - Experte für Compliance, Vertriebsrecht, AGB und E-Commerce.",
            "  - Mehrfach ausgezeichneter Anwalt (Recommended Lawyer in Legal 500 EMEA 2023, 2024, 2025).",
            "- **Kanzlei-Kontaktdaten:**",
            "  - Adresse: Weihburggasse 30/5, A - 1010 Wien",
            "  - Telefon: +43 (1) 4868200-13",
            "  - E-Mail: office@adametz.at",
            "  - Web: www.adametz.at",
            "  - LinkedIn: http://linkedin.com/in/stefan-adametz-1149698"
        ]
    )
    
    # Slide 15: Rechtlicher Disclaimer & Hinweis
    add_slide(
        "Rechtlicher Disclaimer & Hinweis",
        [
            "- **Erstinformation:** Diese Präsentation stellt lediglich eine Erstinformation auf Basis der Gesetzesentwürfe dar.",
            "- **Änderungsvorbehalt:** Die gesetzlichen Details können sich bis zum endgültigen Beschluss des VerbRÄG und WaRUG im Nationalrat noch ändern.",
            "- **Keine Rechtsberatung:** Die Inhalte sind rein informativ und ersetzen keine einfache juristische Beratung.",
            "- **Haftungsausschluss:** Jegliche Haftung des Erstellers der Präsentation oder des Referenten für die Richtigkeit und Vollständigkeit der Inhalte ist ausgeschlossen."
        ]
    )
    
    prs.save(dest_path)
    print(f"Presentation successfully saved to: {dest_path}")

if __name__ == "__main__":
    main()
